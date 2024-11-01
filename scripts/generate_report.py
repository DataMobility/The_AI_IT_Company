
import pandas as pd
import random
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
from io import BytesIO

# Data genereren
def generate_data(num_events=500, power=1.5):
    klanten = ["Ajax", "PSV", "Feyenoord", "FC Utrecht", "AZ Alkmaar", "Heerenveen", "FC Groningen", "Vitesse", "Twente"]
    storingen = [
        ("Power outage", 1 / 4),
        ("Engineer error", 1 / 6),
        ("User error", 1 / 3),
        ("Network congestion", 1 / 5),
        ("Hardware failure", 1 / 10),
        ("Firmware issue", 1 / 8)
    ]
    def random_date():
        start_date = datetime.now() - timedelta(days=365 * 4)
        random_days = random.randint(0, 365 * 4)
        return start_date + timedelta(days=random_days)

    def apply_power(weights, power):
        return [w ** power for w in weights]

    events = []
    storings_namen = [s[0] for s in storingen]
    storings_kansen = [s[1] for s in storingen]
    adjusted_weights = apply_power(storings_kansen, power)

    for _ in range(num_events):
        klant = random.choice(klanten)
        device_type = "Firewall" if random.random() < 0.5 else "Switch"
        device_number = random.randint(1, 4) if device_type == "Firewall" else random.randint(1, 15)
        device = f"{device_type} {device_number}"
        storing = random.choices(storings_namen, weights=adjusted_weights)[0]
        events.append({
            "Datum": random_date().strftime("%Y-%m-%d"),
            "Klant": klant,
            "Device": device,
            "Storing": storing
        })
    return pd.DataFrame(events)

# Data genereren en opslaan
df = generate_data()
df.to_csv("data/historische_storingen.csv", index=False)

# Analyse uitvoeren
storingen_per_type = df['Storing'].value_counts()
total_incidents = len(df)
kans_per_storingstype = (storingen_per_type / total_incidents).to_frame(name="Kans")
verwachte_incidenten_2025 = (kans_per_storingstype * 500).round().astype(int)
analyse = pd.concat([kans_per_storingstype, verwachte_incidenten_2025], axis=1)
analyse.columns = ["Kans", "Verwachte Incidenten 2025"]
analyse.to_csv("output/voorspelde_storingen_2025.csv")

# PowerPoint maken
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Storingsanalyse en Voorspelling voor 2025"
slide.shapes.add_picture("data/D09A8A57-B845-4A47-A351-C9A3897CAAE0.png", Inches(9), Inches(0.5), width=Inches(1.5))
prs.save("output/Storingsanalyse_Voorspelling_2025_Final_Branded_Presentatie.pptx")
