import pandas as pd
import random
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
from io import BytesIO

# Step 1: Generate sample data for incidents
def generate_data(num_events=500, power=1.5):
    clients = ["Ajax", "PSV", "Feyenoord", "FC Utrecht", "AZ Alkmaar", "Heerenveen", "FC Groningen", "Vitesse", "Twente"]
    incidents = [
        ("Power outage", 1 / 4),
        ("Engineer error", 1 / 6),
        ("User error", 1 / 3),
        ("Network congestion", 1 / 5),
        ("Hardware failure", 1 / 10),
        ("Firmware issue", 1 / 8)
    ]
    
    def random_date():
        start_date = datetime.now() - timedelta(days=365 * 4)
        random_days = random.randint(0, 365 * 4)
        return start_date + timedelta(days=random_days)

    def apply_power(weights, power):
        return [w ** power for w in weights]

    events = []
    incident_names = [s[0] for s in incidents]
    incident_weights = [s[1] for s in incidents]
    adjusted_weights = apply_power(incident_weights, power)

    for _ in range(num_events):
        client = random.choice(clients)
        device_type = "Firewall" if random.random() < 0.5 else "Switch"
        device_number = random.randint(1, 4) if device_type == "Firewall" else random.randint(1, 15)
        device = f"{device_type} {device_number}"
        incident = random.choices(incident_names, weights=adjusted_weights)[0]
        events.append({
            "Date": random_date().strftime("%Y-%m-%d"),
            "Client": client,
            "Device": device,
            "Incident": incident
        })
    return pd.DataFrame(events)

# Generate and save data
df = generate_data()
df.to_csv("data/historical_incidents.csv", index=False)

# Perform analysis on generated data and save
incidents_by_type = df['Incident'].value_counts()
total_incidents = len(df)
incident_probabilities = (incidents_by_type / total_incidents).to_frame(name="Probability")
expected_incidents_2025 = (incident_probabilities * 500).round().astype(int)
analysis = pd.concat([incident_probabilities, expected_incidents_2025], axis=1)
analysis.columns = ["Probability", "Expected Incidents 2025"]
analysis.to_csv("output/predicted_incidents_2025.csv")

# Create PowerPoint presentation
prs = Presentation()

# Slide 1: Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Incident Analysis and Prediction for 2025"
subtitle = slide.placeholders[1]
subtitle.text = "BI Report for Football Stadiums by The AI IT Company"

# Apply logo (example logo path)
logo_path = "data/AI_IT_Logo.png"  # Replace with your actual logo path if available
slide.shapes.add_picture(logo_path, Inches(9), Inches(0.5), width=Inches(1.5))

# Slide 2: Historical Incident Overview
# Create chart for incidents per year
df['Year'] = pd.to_datetime(df['Date']).dt.year
incidents_per_year = df['Year'].value_counts().sort_index()

plt.figure(figsize=(8, 6))
plt.bar(incidents_per_year.index, incidents_per_year.values, color='skyblue')
plt.title("Number of Incidents per Year")
plt.xlabel("Year")
plt.ylabel("Number of Incidents")
plt.tight_layout()
buf = BytesIO()
plt.savefig(buf, format='png')
plt.close()

slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Historical Overview of Incidents"
slide.shapes.add_picture(buf, Inches(1), Inches(1.5), width=Inches(8), height=Inches(4.5))

# Slide 3: Incident Types per Client
incident_types_per_client = df.groupby(['Client', 'Incident']).size().unstack(fill_value=0)

plt.figure(figsize=(10, 6))
incident_types_per_client.plot(kind='bar', stacked=True, colormap="tab20")
plt.title("Incident Types per Client")
plt.xlabel("Client")
plt.ylabel("Number of Incidents")
plt.legend(title="Incident Type", bbox_to_anchor=(1.05, 1), loc='upper left')
plt.tight_layout()
buf2 = BytesIO()
plt.savefig(buf2, format='png')
plt.close()

slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Incident Types per Client"
slide.shapes.add_picture(buf2, Inches(1), Inches(1.5), width=Inches(8), height=Inches(4.5))

# Per-Client Analysis
clients_sorted = sorted(df['Client'].unique())
for client in clients_sorted:
    client_data = df[df['Client'] == client]
    years_sorted = sorted(client_data['Year'].unique(), reverse=True)
    for year in years_sorted:
        year_data = client_data[client_data['Year'] == year]
        incidents_per_type = year_data['Incident'].value_counts()

        plt.figure(figsize=(8, 6))
        plt.pie(incidents_per_type, labels=[f"{label} ({count})" for label, count in zip(incidents_per_type.index, incidents_per_type)], autopct='%1.1f%%', startangle=140)
        plt.title(f"Incidents for {client} in {year}")
        plt.tight_layout()
        buf_client = BytesIO()
        plt.savefig(buf_client, format='png')
        plt.close()

        # Add a slide per year per client
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = f"Incidents for {client} in {year}"
        slide.shapes.add_picture(buf_client, Inches(2), Inches(1.5), width=Inches(6), height=Inches(4.5))

# Save the enhanced presentation with detailed slides
prs.save("output/Incident_Analysis_Prediction_2025_Enhanced_Presentation.pptx")
